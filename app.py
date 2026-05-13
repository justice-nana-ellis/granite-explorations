import asyncio
import json
import os
import base64
import html
import io
import logging
import mimetypes
from pathlib import Path
from uuid import uuid4
from time import time
import requests
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import cloudinary
import cloudinary.uploader
import cloudinary.api
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from anthropic import AsyncAnthropic
from fastapi import FastAPI, File, Form, UploadFile, HTTPException, Request
from fastapi.responses import StreamingResponse, HTMLResponse, Response
from pydantic import BaseModel
from typing import Optional, AsyncIterator
from dotenv import load_dotenv

load_dotenv(dotenv_path=Path(__file__).parent / ".env")

cloudinary.config(
    cloud_name=os.getenv("CLOUDINARY_CLOUD_NAME"),
    api_key=os.getenv("CLOUDINARY_API_KEY"),
    api_secret=os.getenv("CLOUDINARY_API_SECRET"),
    secure=True,
)

logging.getLogger("uvicorn.access").setLevel(logging.WARNING)
logging.getLogger("uvicorn.error").setLevel(logging.WARNING)

api_key = os.getenv("CLAUDE_API_KEY")
if not api_key:
    raise ValueError("CLAUDE_API_KEY environment variable not set")

client = AsyncAnthropic(api_key=api_key)
model = os.getenv("CLAUDE_MODEL", "claude-opus-4-7")
chat_model = os.getenv("CLAUDE_CHAT_MODEL", "claude-haiku-4-5-20251001")
port = int(os.getenv("PORT", 8000))

ANALYSIS_SYSTEM_PROMPT = """
You are a world-class data analyst, financial strategist, and forecasting expert — the equivalent of a senior Power BI / Tableau consultant combined with a McKinsey data scientist.

For EVERY response you must produce ALL of the following sections, regardless of how simple the question seems:

1. EXECUTIVE SUMMARY
   A 3–5 sentence high-level summary of the most important findings. Lead with the single most critical insight.

2. KEY INSIGHTS  (minimum 5 bullet points)
   - Use exact numbers, percentages, ratios, and comparisons from the data.
   - Highlight outliers, peaks, troughs, correlations, and unexpected patterns.
   - Be specific — never say "revenue increased" when you can say "revenue grew 23.4% MoM from Jan to Mar".

3. TREND ANALYSIS
   Describe short-term and long-term trends. Identify acceleration or deceleration. Call out seasonality or cyclicality if visible.

4. ANOMALIES & RISKS
   Surface any data anomalies, data quality issues, concerning patterns, or risk signals that a CFO or analyst should be aware of.

5. FORECAST & PROJECTIONS  (always attempt this for numeric data)
   - Extrapolate the next 3–6 periods using the visible trend.
   - State your forecasting methodology (e.g., linear trend, moving average, CAGR).
   - Provide point estimates and a directional confidence range.
   - Example: "Based on a 3-period rolling average, Q3 revenue is forecast at $1.45M (+/- 8%)."

6. STRATEGIC RECOMMENDATIONS  (minimum 3 actionable items)
   Concrete, prioritised actions a decision-maker can take based on the data.

Formatting rules:
- Use **bold headers** for each section.
- Use bullet points — never long unstructured paragraphs.
- Quantify everything: percentages, totals, ratios, MoM/YoY changes.
- Do NOT start with pleasantries like "Great question!" or "Sure!".
- Do NOT say you cannot access the data — the full data summary is already in the message context.
- NEVER say you "cannot create", "cannot display", "cannot show", or "don't have the ability" to make charts, dashboards, or visuals. The visual artifacts are generated automatically by the system — your job is to provide the deep written analysis that accompanies them.
- NEVER suggest the user go to Power BI, Tableau, or any other tool. You ARE the analysis tool.
- NEVER ask clarifying questions. Deliver the full analysis immediately.
"""

app = FastAPI(
    title="Claude Finance API",
    description="Simple Claude API wrapper for financial analysis",
    version="1.0.0",
)

SESSION_STORE = os.getenv("SESSION_STORE", "memory")
SESSION_TTL = int(os.getenv("SESSION_TTL", 3600))
SESSION_MAX_MESSAGES = int(os.getenv("SESSION_MAX_MESSAGES", 10))
SESSION_MAX_COUNT = int(os.getenv("SESSION_MAX_COUNT", 200))

sessions: dict[str, dict] = {}

# ── Cloudinary helpers ────────────────────────────────────────────────

def _cloudinary_delete(public_id: str, resource_type: str = "raw") -> None:
    try:
        cloudinary.uploader.destroy(public_id, resource_type=resource_type)
    except Exception:
        pass


def _delete_all_cloudinary(session: dict) -> None:
    """Delete the uploaded file AND the saved state from Cloudinary."""
    if session.get("cloudinary_id"):
        filename = session.get("file", "")
        rtype = "image" if filename.lower().endswith((".jpg", ".jpeg", ".png", ".gif", ".webp")) else "raw"
        _cloudinary_delete(session["cloudinary_id"], resource_type=rtype)
    if session.get("state_cloudinary_id"):
        _cloudinary_delete(session["state_cloudinary_id"], resource_type="raw")


async def _save_state_to_cloudinary(sid: str, session: dict) -> None:
    """Persist rolling messages + session metadata to Cloudinary as JSON."""
    state = {
        "file": session.get("file"),
        "cloudinary_id": session.get("cloudinary_id"),
        "cloudinary_url": session.get("cloudinary_url"),
        "system": session.get("system"),
        "display": session.get("display", []),
    }
    try:
        raw_json = json.dumps(state).encode()
        result = await asyncio.to_thread(
            cloudinary.uploader.upload,
            io.BytesIO(raw_json),
            public_id=f"sessions/{sid}/state.json",
            resource_type="raw",
            overwrite=True,
        )
        session["state_cloudinary_id"] = result["public_id"]
    except Exception as e:
        print(f"State save failed: {e}")


def _restore_state_from_cloudinary(sid: str, session: dict) -> None:
    """Restore messages + metadata from Cloudinary into an empty session."""
    try:
        cloud_name = os.getenv("CLOUDINARY_CLOUD_NAME")
        url = f"https://res.cloudinary.com/{cloud_name}/raw/upload/sessions/{sid}/state.json"
        resp = requests.get(url, timeout=10)
        if resp.status_code != 200:
            return
        state = resp.json()
        session["file"] = state.get("file")
        session["cloudinary_id"] = state.get("cloudinary_id")
        session["cloudinary_url"] = state.get("cloudinary_url")
        session["system"] = state.get("system") or session["system"]
        session["display"] = state.get("display", [])
        session["state_cloudinary_id"] = f"sessions/{sid}/state.json"
    except Exception:
        pass


# ── Session management ────────────────────────────────────────────────

def _purge_expired() -> None:
    now = time()
    expired = [sid for sid, s in sessions.items() if now - s["last_accessed"] > SESSION_TTL]
    for sid in expired:
        _delete_all_cloudinary(sessions[sid])  # wipe file + state.json from Cloudinary
        del sessions[sid]


def _get_or_create_session(sid: str, system: str) -> dict:
    _purge_expired()
    if sid in sessions:
        sessions[sid]["last_accessed"] = time()
        return sessions[sid]

    if len(sessions) >= SESSION_MAX_COUNT:
        oldest = min(sessions, key=lambda k: sessions[k]["last_accessed"])
        _delete_all_cloudinary(sessions[oldest])  # eviction: wipe Cloudinary too
        del sessions[oldest]

    sessions[sid] = {
        "system": system,
        "messages": [],
        "display": [],
        "file": None,
        "cloudinary_id": None,
        "cloudinary_url": None,
        "state_cloudinary_id": None,
        "df": None,
        "df_summary": None,
        "last_accessed": time(),
    }

    # Try to restore from Cloudinary (server restart / TTL expiry)
    _restore_state_from_cloudinary(sid, sessions[sid])
    return sessions[sid]


def _reload_df_from_cloudinary(session: dict) -> None:
    if not session.get("cloudinary_url") or not session.get("file"):
        return
    try:
        resp = requests.get(session["cloudinary_url"], timeout=30)
        resp.raise_for_status()
        raw = resp.content
        filename: str = session["file"]
        if filename.endswith(".csv"):
            session["df"] = pd.read_csv(io.BytesIO(raw), low_memory=False)
        elif filename.endswith((".xlsx", ".xls")):
            session["df"] = pd.read_excel(io.BytesIO(raw), engine="openpyxl")
        if session["df"] is not None:
            session["df_summary"] = _df_summary(session["df"], filename)
    except Exception:
        pass


def _session_meta(sid: str) -> dict:
    s = sessions[sid]
    return {"session_id": sid, "file": s["file"], "messages": s["display"]}


def _trim_messages(session: dict) -> None:
    if len(session["messages"]) >= SESSION_MAX_MESSAGES:
        first = session["messages"][:1]
        rest = session["messages"][1:][-(SESSION_MAX_MESSAGES - 2):]
        session["messages"] = first + rest
    if len(session["display"]) >= SESSION_MAX_MESSAGES:
        session["display"] = session["display"][-SESSION_MAX_MESSAGES:]


# ── Models ────────────────────────────────────────────────────────────

class ChatRequest(BaseModel):
    message: str
    system_prompt: Optional[str] = None
    session_id: Optional[str] = None
    response_format: Optional[str] = None


class ChatResponse(BaseModel):
    reply: str
    model: str


# ── Startup ───────────────────────────────────────────────────────────

@app.on_event("startup")
async def startup():
    print(f"✓ Analysis model : {model}")
    print(f"✓ Chat model     : {chat_model}")
    print(f"✓ Server running on http://0.0.0.0:{port}")


@app.get("/")
def root():
    return {"status": "ok", "model": model, "port": port}


@app.get("/health")
def health():
    return {"status": "healthy"}


# ── Streaming helper ──────────────────────────────────────────────────

async def _stream_and_save(session_id: str, use_model: str, max_tokens: int) -> AsyncIterator[str]:
    session = sessions[session_id]
    full_reply = []

    async with client.messages.stream(
        model=use_model,
        max_tokens=max_tokens,
        system=session["system"],
        messages=session["messages"],
    ) as stream:
        async for text in stream.text_stream:
            full_reply.append(text)
            yield f"data: {text}\n\n"

    reply_text = "".join(full_reply)
    session["messages"].append({"role": "assistant", "content": reply_text})
    session["display"].append({"role": "assistant", "content": reply_text})

    # Save updated messages to Cloudinary in the background
    asyncio.create_task(_save_state_to_cloudinary(session_id, session))


# ── Chat endpoint ─────────────────────────────────────────────────────

@app.post("/chat")
async def chat(request: ChatRequest):
    sid = request.session_id or str(uuid4())
    session = _get_or_create_session(
        sid, request.system_prompt or "You are a helpful financial assistant."
    )

    _trim_messages(session)

    if session.get("df") is None and session.get("cloudinary_url"):
        await asyncio.to_thread(_reload_df_from_cloudinary, session)

    if session.get("df_summary"):
        message_with_context = f"[File context]\n{session['df_summary']}\n\n[Question]\n{request.message}"
    else:
        message_with_context = request.message

    session["messages"].append({"role": "user", "content": message_with_context})
    session["display"].append({"role": "user", "content": request.message})

    async def _generate():
        async for chunk in _stream_and_save(sid, chat_model, 1024):
            yield chunk
        sessions[sid]["last_accessed"] = time()
        yield f"event: meta\ndata: {json.dumps(_session_meta(sid))}\n\n"

    return StreamingResponse(_generate(), media_type="text/event-stream")


# ── File upload helpers ───────────────────────────────────────────────

EXCEL_TYPES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
}

WORD_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}

PPTX_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}


def _extract_docx(raw: bytes) -> str:
    doc = DocxDocument(io.BytesIO(raw))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(raw: bytes) -> str:
    prs = Presentation(io.BytesIO(raw))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _df_summary(df: pd.DataFrame, filename: str) -> str:
    numeric_cols = df.select_dtypes(include="number").columns.tolist()
    date_cols = [c for c in df.columns if "date" in c.lower()]
    cat_cols = [c for c in df.select_dtypes(exclude="number").columns if c not in date_cols]

    lines = [
        f"File: {filename}",
        f"Rows: {len(df):,} | Columns: {len(df.columns)}",
        f"Columns: {', '.join(df.columns.tolist())}",
    ]
    if date_cols:
        for dc in date_cols:
            uniq = df[dc].dropna().unique()
            lines.append(f"Unique {dc} values ({len(uniq)}): {', '.join(sorted(str(v) for v in uniq))}")
    if numeric_cols:
        lines.append("\nNumeric column totals (all rows):")
        for col in numeric_cols:
            col_data = df[col].dropna()
            lines.append(
                f"  {col}: total={col_data.sum():,.4f}, mean={col_data.mean():,.4f}, "
                f"min={col_data.min():,.4f}, max={col_data.max():,.4f}, count={len(col_data):,}"
            )
    if date_cols and numeric_cols:
        primary_date = date_cols[0]
        lines.append(f"\nNumeric totals grouped by {primary_date}:")
        grouped = df.groupby(primary_date)[numeric_cols].sum()
        for date_val, row in grouped.iterrows():
            row_parts = ", ".join(f"{col}={row[col]:,.2f}" for col in numeric_cols)
            lines.append(f"  {date_val}: {row_parts}")
    if cat_cols:
        lines.append("\nCategorical column summaries:")
        for col in cat_cols[:10]:
            uniq = df[col].dropna().unique()
            sample = ", ".join(str(v) for v in uniq[:10])
            lines.append(f"  {col}: {len(uniq)} unique values — e.g. {sample}")
    lines.append(f"\nFirst 5 rows:\n{df.head(5).to_string(index=False)}")
    return "\n".join(lines)


# ── Upload endpoint ───────────────────────────────────────────────────

@app.post("/upload")
async def upload_and_ask(
    file: UploadFile = File(...),
    question: str = Form(...),
    system_prompt: Optional[str] = Form(None),
    session_id: Optional[str] = Form(None),
):
    content_type = file.content_type or ""
    if not content_type or "/" not in content_type:
        content_type, _ = mimetypes.guess_type(file.filename or "")
        content_type = content_type or "application/octet-stream"
    if content_type == "image/jpg":
        content_type = "image/jpeg"
    raw = await file.read()

    if not raw:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    parsed_df = None
    summary = None

    if content_type == "text/csv":
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="CSV file could not be decoded as UTF-8")
        parsed_df = await asyncio.to_thread(pd.read_csv, io.StringIO(text), low_memory=False)
        summary = await asyncio.to_thread(_df_summary, parsed_df, file.filename)
        user_content = [{"type": "text", "text": f"{summary}\n\nQuestion: {question}"}]

    elif content_type in EXCEL_TYPES:
        try:
            parsed_df = await asyncio.to_thread(pd.read_excel, io.BytesIO(raw), engine="openpyxl")
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Could not read Excel file: {e}")
        summary = await asyncio.to_thread(_df_summary, parsed_df, file.filename)
        user_content = [{"type": "text", "text": f"{summary}\n\nQuestion: {question}"}]

    elif content_type in WORD_TYPES:
        try:
            text_content = await asyncio.to_thread(_extract_docx, raw)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Could not read Word file: {e}")
        user_content = [{"type": "text", "text": f"File: {file.filename}\n\n{text_content}\n\n{question}"}]

    elif content_type in PPTX_TYPES:
        try:
            text_content = await asyncio.to_thread(_extract_pptx, raw)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Could not read PowerPoint file: {e}")
        user_content = [{"type": "text", "text": f"File: {file.filename}\n\n{text_content}\n\n{question}"}]

    elif content_type.startswith("text/"):
        try:
            text_content = raw.decode("utf-8")
        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="File could not be decoded as UTF-8 text")
        user_content = [{"type": "text", "text": f"File: {file.filename}\n\n{text_content}\n\n{question}"}]

    elif content_type in ("image/jpeg", "image/png", "image/gif", "image/webp"):
        encoded = base64.standard_b64encode(raw).decode("utf-8")
        user_content = [
            {"type": "image", "source": {"type": "base64", "media_type": content_type, "data": encoded}},
            {"type": "text", "text": question},
        ]

    elif content_type == "application/pdf":
        encoded = base64.standard_b64encode(raw).decode("utf-8")
        user_content = [
            {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": encoded}},
            {"type": "text", "text": question},
        ]

    else:
        raise HTTPException(
            status_code=415,
            detail=f"Unsupported file type '{content_type}'. Supported: CSV, Excel, images (jpeg/png/gif/webp), PDF, and plain text.",
        )

    sid = session_id or str(uuid4())
    session = _get_or_create_session(
        sid, system_prompt or "You are a helpful assistant. Analyse the provided file and answer the user's question precisely."
    )

    _trim_messages(session)

    session["file"] = file.filename
    if parsed_df is not None:
        session["df"] = parsed_df
        session["df_summary"] = summary

    session["messages"].append({"role": "user", "content": user_content})
    session["display"].append({"role": "user", "content": question})

    # Upload file + initial state to Cloudinary in the background
    async def _background_upload():
        resource_type = "image" if content_type.startswith("image/") else "raw"
        try:
            result = await asyncio.to_thread(
                cloudinary.uploader.upload,
                io.BytesIO(raw),
                public_id=f"sessions/{sid}/{file.filename}",
                resource_type=resource_type,
                overwrite=True,
            )
            session["cloudinary_id"] = result["public_id"]
            session["cloudinary_url"] = result["secure_url"]
        except Exception as e:
            print(f"Cloudinary file upload failed: {e}")
        # Save initial state (messages so far) once file upload is done
        await _save_state_to_cloudinary(sid, session)

    asyncio.create_task(_background_upload())

    async def _generate():
        async for chunk in _stream_and_save(sid, model, 2048):
            yield chunk
        sessions[sid]["last_accessed"] = time()
        yield f"event: meta\ndata: {json.dumps(_session_meta(sid))}\n\n"

    return StreamingResponse(_generate(), media_type="text/event-stream")


# ── Session list ─────────────────────────────────────────────────────

@app.get("/sessions")
def list_sessions():
    _purge_expired()
    return {
        "count": len(sessions),
        "sessions": [_session_meta(sid) for sid in sessions],
    }


# ── Files list ────────────────────────────────────────────────────────

@app.get("/files")
async def list_files():
    """List all files stored in Cloudinary under the sessions/ folder."""
    try:
        def _fetch_all():
            results = []
            for rtype in ("image", "raw"):
                next_cursor = None
                while True:
                    kwargs = dict(
                        resource_type=rtype,
                        type="upload",
                        prefix="sessions/",
                        max_results=500,
                    )
                    if next_cursor:
                        kwargs["next_cursor"] = next_cursor
                    resp = cloudinary.api.resources(**kwargs)
                    results.extend(resp.get("resources", []))
                    next_cursor = resp.get("next_cursor")
                    if not next_cursor:
                        break
            return results

        all_resources = await asyncio.to_thread(_fetch_all)
        files = [
            {
                "public_id": r["public_id"],
                "url": r["secure_url"],
                "bytes": r["bytes"],
                "format": r.get("format") or os.path.splitext(r["public_id"])[-1].lstrip(".") or "unknown",
                "resource_type": r.get("resource_type"),
                "created_at": r.get("created_at"),
            }
            for r in all_resources
            if not r["public_id"].endswith("state.json")
        ]
        return {"count": len(files), "files": files}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Cloudinary error: {e}")


# ── Session delete ────────────────────────────────────────────────────

@app.delete("/session/{session_id}")
def clear_session(session_id: str):
    session = sessions.pop(session_id, None)
    if session:
        _delete_all_cloudinary(session)
    return {"cleared": session_id}


# ── New helper utilities ─────────────────────────────────────────────

def _normalized_content_type(upload: UploadFile) -> str:
    content_type = upload.content_type or ""
    if not content_type or "/" not in content_type or content_type == "application/octet-stream":
        guessed_type, _ = mimetypes.guess_type(upload.filename or "")
        content_type = guessed_type or "application/octet-stream"
    if content_type == "image/jpg":
        content_type = "image/jpeg"
    return content_type


def _wants_visual_render(message: str) -> bool:
    text = message.lower()
    triggers = ("chart", "plot", "graph", "visual", "dashboard", "power bi", "render", "forecast", "projection")
    return any(trigger in text for trigger in triggers)


def _extract_assistant_text(response) -> str:
    text_parts = []
    for block in getattr(response, "content", []) or []:
        maybe_text = getattr(block, "text", None)
        if maybe_text:
            text_parts.append(maybe_text)
    return "\n".join(text_parts).strip() or "No assistant response generated."


async def _ensure_session_df(session: dict) -> pd.DataFrame:
    if session.get("df") is None and session.get("cloudinary_url"):
        await asyncio.to_thread(_reload_df_from_cloudinary, session)
    df = session.get("df")
    if df is None:
        raise HTTPException(status_code=400, detail="No tabular data found for this session. Upload CSV/Excel first.")
    return df


async def _build_chat_visual_bundle(sid: str, session: dict, base_url: str) -> dict:
    artifacts = await _build_visual_artifacts(sid, session, chart_type="auto", max_charts=3)
    forecast = await _build_forecast_artifact(sid, session)
    return {
        "session_id": sid,
        "file": session.get("file"),
        "source_file_url": session.get("cloudinary_url"),
        "dashboard_url": f"{base_url}/dashboard/{sid}/render",
        "artifacts_url": f"{base_url}/artifacts/{sid}",
        "forecast_url": f"{base_url}/forecast/{sid}",
        "artifacts": artifacts,
        "forecast": forecast,
    }


def _make_chart_bytes(df: pd.DataFrame, chart_type: str, x_col: str, y_col: str) -> bytes:
    fig, ax = plt.subplots(figsize=(11, 6))
    working = df[[x_col, y_col]].dropna().copy()

    if chart_type in ("auto", "bar"):
        is_date_like = any(k in x_col.lower() for k in ("date", "month", "time", "year"))
        chosen = "line" if is_date_like else "bar"
    else:
        chosen = chart_type

    if chosen == "line":
        working.plot(x=x_col, y=y_col, kind="line", marker="o", ax=ax, legend=False, color="#2563eb")
        ax.fill_between(range(len(working)), working[y_col], alpha=0.08, color="#2563eb")
    elif chosen == "scatter":
        working.plot(x=x_col, y=y_col, kind="scatter", ax=ax, color="#7c3aed")
    else:
        working.head(50).plot(x=x_col, y=y_col, kind="bar", ax=ax, legend=False, color="#2563eb", edgecolor="#1d4ed8")

    ax.set_title(f"{y_col} by {x_col}", fontsize=14, fontweight="bold", pad=12)
    ax.set_xlabel(x_col, fontsize=11)
    ax.set_ylabel(y_col, fontsize=11)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", linestyle="--", alpha=0.4)
    fig.patch.set_facecolor("#ffffff")
    ax.set_facecolor("#fafafa")
    fig.tight_layout()

    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", dpi=140)
    plt.close(fig)
    buffer.seek(0)
    return buffer.getvalue()


def _resolve_chart_specs(df: pd.DataFrame, chart_type: str, max_charts: int) -> list[tuple[str, str, str]]:
    numeric_cols = df.select_dtypes(include="number").columns.tolist()
    all_cols = df.columns.tolist()
    if not numeric_cols or len(all_cols) < 2:
        raise HTTPException(status_code=400, detail="Not enough usable columns to create charts.")

    x_candidates = [column for column in all_cols if column not in numeric_cols]
    x_col = x_candidates[0] if x_candidates else all_cols[0]
    chosen_chart = "bar" if chart_type == "auto" else chart_type

    specs = []
    for y_col in numeric_cols:
        if y_col == x_col:
            continue
        specs.append((chosen_chart, x_col, y_col))
        if len(specs) >= max_charts:
            break

    if not specs:
        raise HTTPException(status_code=400, detail="Could not determine chart columns from uploaded data.")
    return specs


async def _build_visual_artifacts(sid: str, session: dict, chart_type: str = "auto", max_charts: int = 3) -> list[dict]:
    df = await _ensure_session_df(session)
    specs = _resolve_chart_specs(df, chart_type=chart_type, max_charts=max_charts)
    artifacts = []

    for chosen_chart, x_col, y_col in specs:
        chart_bytes = await asyncio.to_thread(_make_chart_bytes, df, chosen_chart, x_col, y_col)
        upload = await asyncio.to_thread(
            cloudinary.uploader.upload,
            io.BytesIO(chart_bytes),
            public_id=f"sessions/{sid}/artifacts/{chosen_chart}_{y_col}",
            resource_type="image",
            overwrite=True,
        )
        artifacts.append({
            "title": f"{chosen_chart.title()} - {y_col} by {x_col}",
            "chart_type": chosen_chart,
            "x": x_col,
            "y": y_col,
            "public_id": upload["public_id"],
            "url": upload["secure_url"],
        })

    return artifacts


def _find_forecast_columns(df: pd.DataFrame) -> tuple[str, str, pd.DataFrame]:
    working = df.copy()
    date_candidates = [c for c in working.columns if any(k in c.lower() for k in ("date", "month", "time"))]
    if not date_candidates:
        raise HTTPException(status_code=400, detail="Could not find a date-like column for forecasting.")

    date_col = date_candidates[0]
    working[date_col] = pd.to_datetime(working[date_col], errors="coerce")
    working = working.dropna(subset=[date_col])

    numeric_cols = [c for c in working.select_dtypes(include="number").columns.tolist() if c != date_col]
    if not numeric_cols:
        raise HTTPException(status_code=400, detail="Could not find a numeric column for forecasting.")

    value_col = numeric_cols[0]
    grouped = working.groupby(date_col, as_index=False)[value_col].sum().sort_values(date_col)
    if len(grouped) < 2:
        raise HTTPException(status_code=400, detail="Need at least two dated data points to build a forecast.")
    return date_col, value_col, grouped


def _make_forecast_chart_bytes(actual_df: pd.DataFrame, forecast_df: pd.DataFrame, date_col: str, value_col: str) -> bytes:
    fig, ax = plt.subplots(figsize=(11, 6))
    actual_df.plot(x=date_col, y=value_col, kind="line", marker="o", ax=ax, label="Actual")
    forecast_df.plot(x=date_col, y=value_col, kind="line", marker="o", linestyle="--", ax=ax, label="Forecast", color="#d97706")
    ax.set_title(f"Forecast for {value_col}")
    ax.set_xlabel(date_col)
    ax.set_ylabel(value_col)
    ax.legend()
    fig.tight_layout()

    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", dpi=130)
    plt.close(fig)
    buffer.seek(0)
    return buffer.getvalue()


async def _build_forecast_artifact(sid: str, session: dict, periods: int = 3) -> dict:
    df = await _ensure_session_df(session)
    date_col, value_col, grouped = _find_forecast_columns(df)

    recent = grouped.tail(min(4, len(grouped))).reset_index(drop=True)
    slope = (recent[value_col].iloc[-1] - recent[value_col].iloc[0]) / max(len(recent) - 1, 1)
    last_date = recent[date_col].iloc[-1]
    step = recent[date_col].diff().dropna().median()
    if pd.isna(step) or step <= pd.Timedelta(0):
        step = pd.Timedelta(days=30)

    forecast_rows = []
    last_value = recent[value_col].iloc[-1]
    for index in range(1, periods + 1):
        forecast_rows.append({date_col: last_date + (step * index), value_col: float(last_value + (slope * index))})

    forecast_df = pd.DataFrame(forecast_rows)
    chart_bytes = await asyncio.to_thread(_make_forecast_chart_bytes, recent, forecast_df, date_col, value_col)
    upload = await asyncio.to_thread(
        cloudinary.uploader.upload,
        io.BytesIO(chart_bytes),
        public_id=f"sessions/{sid}/artifacts/forecast_{value_col}",
        resource_type="image",
        overwrite=True,
    )

    serialized_rows = [
        {date_col: pd.Timestamp(row[date_col]).strftime("%Y-%m-%d"), value_col: row[value_col]}
        for row in forecast_rows
    ]

    return {
        "title": f"Forecast - {value_col}",
        "chart_type": "forecast",
        "x": date_col,
        "y": value_col,
        "public_id": upload["public_id"],
        "url": upload["secure_url"],
        "forecast_rows": serialized_rows,
        "summary": f"Forecasted {value_col} for the next {periods} period(s) based on the recent trend.",
    }


def _build_kpi_cards_html(session: dict) -> str:
    df: pd.DataFrame = session.get("df")
    if df is None:
        return ""
    numeric_cols = df.select_dtypes(include="number").columns.tolist()[:6]
    cards = ""
    palette = ["#2563eb", "#16a34a", "#d97706", "#7c3aed", "#dc2626", "#0891b2"]
    for i, col in enumerate(numeric_cols):
        data = df[col].dropna()
        total = data.sum()
        mean = data.mean()
        color = palette[i % len(palette)]
        def _fmt(v):
            if abs(v) >= 1_000_000:
                return f"{v/1_000_000:,.2f}M"
            if abs(v) >= 1_000:
                return f"{v/1_000:,.1f}K"
            return f"{v:,.2f}"
        cards += (
            f"<div style='background:#fff;border-top:4px solid {color};border-radius:14px;padding:16px 18px;"
            f"box-shadow:0 1px 6px rgba(0,0,0,.06);min-width:140px'>"
            f"<div style='font:600 11px Arial,sans-serif;color:#6b7280;text-transform:uppercase;letter-spacing:.6px;margin-bottom:6px'>{html.escape(col)}</div>"
            f"<div style='font:700 24px Arial,sans-serif;color:#111827;line-height:1.1'>{_fmt(total)}</div>"
            f"<div style='font:500 12px Arial,sans-serif;color:#6b7280;margin-top:4px'>avg {_fmt(mean)}</div>"
            f"</div>"
        )
    if not cards:
        return ""
    return f"<div style='display:flex;flex-wrap:wrap;gap:14px'>{cards}</div>"


def _build_chartjs_html(session: dict):
    df: pd.DataFrame = session.get("df")
    if df is None:
        return ("", "")
    all_cols = df.columns.tolist()
    numeric_cols = df.select_dtypes(include="number").columns.tolist()
    x_candidates = [c for c in all_cols if c not in numeric_cols]
    if not numeric_cols or not x_candidates:
        return ("", "")
    x_col = x_candidates[0]
    is_ts = any(k in x_col.lower() for k in ("date", "month", "time", "year", "period"))
    chart_cols = numeric_cols[:6]
    labels_json = json.dumps([str(v) for v in df[x_col].fillna("").tolist()[:100]])
    palette = ["#2563eb", "#16a34a", "#d97706", "#7c3aed", "#dc2626", "#0891b2"]
    chart_type_js = "line" if is_ts else "bar"
    inits = []

    overview_datasets = []
    for i, col in enumerate(chart_cols):
        vals_json = json.dumps([round(float(v), 4) if pd.notna(v) else None for v in df[col].tolist()[:100]])
        color = palette[i % len(palette)]
        fill_val = "true" if (is_ts and i == 0) else "false"
        pt_r = 3 if len(df) <= 40 else 1
        overview_datasets.append(
            "{label:" + json.dumps(col) + ",data:" + vals_json + ","
            "borderColor:'" + color + "',backgroundColor:'" + color + "33',"
            "fill:" + fill_val + ",tension:0.35,pointRadius:" + str(pt_r) + ",borderWidth:2.5}"
        )
    inits.append(
        "(function(){var el=document.getElementById('chartOverview');"
        "if(!el)return;if(el._ci){el._ci.resize();return;}"
        "el._ci=new Chart(el.getContext('2d'),{type:'" + chart_type_js + "',"
        "data:{labels:" + labels_json + ",datasets:[" + ",".join(overview_datasets) + "]},"
        "options:{responsive:true,maintainAspectRatio:false,"
        "interaction:{mode:'index',intersect:false},"
        "plugins:{legend:{position:'top',labels:{font:{size:12},padding:14}},"
        "tooltip:{callbacks:{label:function(c){return ' '+c.dataset.label+': '+Number(c.raw).toLocaleString();}}}},"
        "scales:{x:{ticks:{maxRotation:45,font:{size:11},maxTicksLimit:24}},"
        "y:{ticks:{callback:function(v){return v>=1e6?(v/1e6).toFixed(1)+'M':v>=1e3?(v/1e3).toFixed(1)+'K':v;},font:{size:11}}}}}});"
        "if(!window._chartInstances)window._chartInstances=[];"
        "window._chartInstances.push(el._ci);})();"
    )

    totals = [round(float(df[c].sum()), 2) for c in chart_cols]
    totals_json = json.dumps(totals)
    donut_labels_json = json.dumps(chart_cols)
    donut_colors_json = json.dumps(palette[:len(chart_cols)])
    inits.append(
        "(function(){var el=document.getElementById('chartDonut');"
        "if(!el)return;if(el._ci){el._ci.resize();return;}"
        "el._ci=new Chart(el.getContext('2d'),{type:'doughnut',"
        "data:{labels:" + donut_labels_json + ",datasets:[{data:" + totals_json + ","
        "backgroundColor:" + donut_colors_json + ",borderWidth:3,borderColor:'#fff',hoverOffset:10}]},"
        "options:{responsive:true,maintainAspectRatio:false,cutout:'62%',"
        "plugins:{legend:{position:'right',labels:{font:{size:12},padding:14}},"
        "tooltip:{callbacks:{label:function(c){"
        "var t=c.chart.data.datasets[0].data.reduce(function(a,b){return a+b;},0);"
        "return ' '+c.label+': '+Number(c.raw).toLocaleString()+' ('+(c.raw/t*100).toFixed(1)+'%)';}}}}}});"
        "if(!window._chartInstances)window._chartInstances=[];"
        "window._chartInstances.push(el._ci);})();"
    )

    last_row = df[chart_cols].dropna(how="all").tail(1)
    if len(last_row) > 0:
        snap_vals = [round(float(last_row[c].iloc[0]), 2) if pd.notna(last_row[c].iloc[0]) else 0 for c in chart_cols]
    else:
        snap_vals = [round(float(df[c].mean()), 2) for c in chart_cols]
    snap_vals_json = json.dumps(snap_vals)
    hbar_colors_json = json.dumps([p + "cc" for p in palette[:len(chart_cols)]])
    inits.append(
        "(function(){var el=document.getElementById('chartHBar');"
        "if(!el)return;if(el._ci){el._ci.resize();return;}"
        "el._ci=new Chart(el.getContext('2d'),{type:'bar',"
        "data:{labels:" + donut_labels_json + ",datasets:[{label:'Latest snapshot',"
        "data:" + snap_vals_json + ",backgroundColor:" + hbar_colors_json + ",borderRadius:6,borderWidth:0}]},"
        "options:{responsive:true,maintainAspectRatio:false,indexAxis:'y',"
        "plugins:{legend:{display:false},"
        "tooltip:{callbacks:{label:function(c){return ' '+Number(c.raw).toLocaleString();}}}},"
        "scales:{x:{ticks:{callback:function(v){return v>=1e6?(v/1e6).toFixed(1)+'M':v>=1e3?(v/1e3).toFixed(1)+'K':v;},font:{size:11}}},"
        "y:{ticks:{font:{size:12}}}}}});"
        "if(!window._chartInstances)window._chartInstances=[];"
        "window._chartInstances.push(el._ci);})();"
    )

    if is_ts and len(chart_cols) >= 2:
        stacked_datasets = []
        for i, col in enumerate(chart_cols[:4]):
            vals_json = json.dumps([round(float(v), 4) if pd.notna(v) else 0 for v in df[col].tolist()[:100]])
            color = palette[i % len(palette)]
            stacked_datasets.append(
                "{label:" + json.dumps(col) + ",data:" + vals_json + ","
                "backgroundColor:'" + color + "99',borderColor:'" + color + "',"
                "borderWidth:1,borderRadius:2,stack:'s1'}"
            )
        inits.append(
            "(function(){var el=document.getElementById('chartStacked');"
            "if(!el)return;if(el._ci){el._ci.resize();return;}"
            "el._ci=new Chart(el.getContext('2d'),{type:'bar',"
            "data:{labels:" + labels_json + ",datasets:[" + ",".join(stacked_datasets) + "]},"
            "options:{responsive:true,maintainAspectRatio:false,"
            "interaction:{mode:'index',intersect:false},"
            "plugins:{legend:{position:'top',labels:{font:{size:12}}},"
            "tooltip:{callbacks:{label:function(c){return ' '+c.dataset.label+': '+Number(c.raw).toLocaleString();}}}},"
            "scales:{x:{stacked:true,ticks:{maxRotation:45,font:{size:11},maxTicksLimit:24}},"
            "y:{stacked:true,ticks:{callback:function(v){return v>=1e6?(v/1e6).toFixed(1)+'M':v>=1e3?(v/1e3).toFixed(1)+'K':v;},font:{size:11}}}}}});"
            "if(!window._chartInstances)window._chartInstances=[];"
            "window._chartInstances.push(el._ci);})();"
        )
        stacked_html = (
            "<div style='background:#fff;border:1px solid #e5e7eb;border-radius:14px;padding:16px;margin-bottom:18px'>"
            "<div style='font:700 13px Arial,sans-serif;color:#1e3a8a;margin-bottom:10px'>&#9641; Stacked Period Comparison</div>"
            "<div style='position:relative;height:300px'><canvas id='chartStacked'></canvas></div>"
            "</div>"
        )
    else:
        stacked_html = ""

    spark_canvases = []
    for i, col in enumerate(chart_cols):
        cid = f"chartC{i}"
        vals_json = json.dumps([round(float(v), 4) if pd.notna(v) else None for v in df[col].tolist()[:100]])
        color = palette[i % len(palette)]
        fill_val = "true" if (is_ts and i == 0) else "false"
        pt_r = 3 if len(df) <= 40 else 1
        spark_canvases.append(
            f"<div style='background:#fafafa;border:1px solid #e5e7eb;border-radius:14px;padding:14px'>"
            f"<div style='font:700 12px Arial,sans-serif;color:#374151;margin-bottom:8px'><strong>{html.escape(col)}</strong></div>"
            f"<div style='position:relative;height:190px'><canvas id='{cid}'></canvas></div>"
            f"</div>"
        )
        inits.append(
            "(function(){var el=document.getElementById('" + cid + "');"
            "if(!el)return;if(el._ci){el._ci.resize();return;}"
            "el._ci=new Chart(el.getContext('2d'),{type:'" + chart_type_js + "',"
            "data:{labels:" + labels_json + ",datasets:[{label:" + json.dumps(col) + ",data:" + vals_json + ","
            "borderColor:'" + color + "',backgroundColor:'" + color + "22',"
            "fill:" + fill_val + ",tension:0.35,pointRadius:" + str(pt_r) + ",borderWidth:2}]},"
            "options:{responsive:true,maintainAspectRatio:false,"
            "plugins:{legend:{display:false},"
            "tooltip:{callbacks:{label:function(c){return c.dataset.label+': '+Number(c.raw).toLocaleString();}}}},"
            "scales:{x:{ticks:{maxRotation:45,font:{size:10},maxTicksLimit:16}},"
            "y:{ticks:{callback:function(v){return v>=1e6?(v/1e6).toFixed(1)+'M':v>=1e3?(v/1e3).toFixed(1)+'K':v;},font:{size:10}}}}}});"
            "if(!window._chartInstances)window._chartInstances=[];"
            "window._chartInstances.push(el._ci);})();"
        )

    overview_block = (
        "<div style='background:#fff;border:1px solid #e5e7eb;border-radius:14px;padding:16px;margin-bottom:18px'>"
        "<div style='font:700 14px Arial,sans-serif;color:#1e3a8a;margin-bottom:12px'>&#128200; All Metrics Overview</div>"
        "<div style='position:relative;height:340px'><canvas id='chartOverview'></canvas></div>"
        "</div>"
    )
    two_col_block = (
        "<div style='display:grid;grid-template-columns:1fr 1fr;gap:16px;margin-bottom:18px'>"
        "<div style='background:#fff;border:1px solid #e5e7eb;border-radius:14px;padding:16px'>"
        "<div style='font:700 13px Arial,sans-serif;color:#1e3a8a;margin-bottom:10px'>&#11835; Total Distribution</div>"
        "<div style='position:relative;height:260px'><canvas id='chartDonut'></canvas></div>"
        "</div>"
        "<div style='background:#fff;border:1px solid #e5e7eb;border-radius:14px;padding:16px'>"
        "<div style='font:700 13px Arial,sans-serif;color:#1e3a8a;margin-bottom:10px'>&#128293; Metric Snapshot</div>"
        "<div style='position:relative;height:260px'><canvas id='chartHBar'></canvas></div>"
        "</div>"
        "</div>"
    )
    spark_grid_block = (
        "<div style='font:700 14px Arial,sans-serif;color:#1e3a8a;margin:4px 0 12px'>&#9889; Individual Trends</div>"
        "<div style='display:grid;grid-template-columns:repeat(auto-fit,minmax(300px,1fr));gap:14px;margin-bottom:4px'>"
        + "".join(spark_canvases)
        + "</div>"
    )

    init_js = "window._initIC=function(){"
    init_js += "if(typeof Chart==='undefined'){setTimeout(window._initIC,200);return;}"
    init_js += "".join(inits)
    init_js += "};"
    init_js += (
        "if(document.readyState==='loading'){"
        "document.addEventListener('DOMContentLoaded',function(){setTimeout(window._initIC,50)});"
        "}else{setTimeout(window._initIC,50);}"
        "window.addEventListener('load',function(){setTimeout(window._initIC,50);});"
    )

    return (overview_block + two_col_block + stacked_html + spark_grid_block, init_js)


def _build_data_table_html(session: dict, max_rows: int = 10) -> str:
    df: pd.DataFrame = session.get("df")
    if df is None:
        return ""
    preview = df.head(max_rows)
    th_cells = "".join(f"<th>{html.escape(str(c))}</th>" for c in preview.columns)
    rows_html = ""
    for _, row in preview.iterrows():
        cells = "".join(f"<td>{html.escape(str(v))}</td>" for v in row)
        rows_html += f"<tr>{cells}</tr>"
    return (
        f"<div style='background:#fff;border:1px solid #e5e7eb;border-radius:16px;padding:20px;"
        f"box-shadow:0 1px 6px rgba(0,0,0,.05);overflow-x:auto'>"
        f"<div style='font:700 15px Arial,sans-serif;color:#1e3a8a;margin-bottom:12px'>&#128196; Data Preview (first {len(preview)} rows)</div>"
        f"<table style='width:100%;border-collapse:collapse;font-size:13px;white-space:nowrap'>"
        f"<thead><tr style='background:#1e3a8a;color:#fff'>{th_cells}</tr></thead>"
        f"<tbody>{rows_html}</tbody>"
        f"</table>"
        f"<div style='font:500 12px Arial,sans-serif;color:#9ca3af;margin-top:8px'>{len(df):,} total rows × {len(df.columns)} columns</div>"
        f"</div>"
    )


def _analysis_to_html(analysis: str) -> str:
    import re
    result = html.escape(analysis)
    result = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", result)
    result = re.sub(r"(?m)^#{2,3}\s+(.+)$", r"<h3 style='font-size:15px;font-weight:700;color:#1e3a8a;margin:16px 0 6px'>\1</h3>", result)
    result = re.sub(
        r"(?m)^(\d+\.\s+)([A-Z &]+(?:\s+&amp;\s+[A-Z]+)*)$",
        r"<h3 style='font-size:15px;font-weight:700;color:#1e3a8a;margin:18px 0 6px'>\1\2</h3>",
        result
    )
    result = re.sub(r"(?m)^\s*[-\u2022]\s+(.+)$", r"<li style='margin:4px 0 4px 18px'>\1</li>", result)
    result = re.sub(r"(<li[^>]*>.*?</li>(?:\s*<li[^>]*>.*?</li>)*)", r"<ul style='list-style:disc;padding:0;margin:6px 0'>\1</ul>", result, flags=re.DOTALL)
    result = result.replace("\n", "<br>")
    return result


def _build_jsx_snippet(session: dict, analysis: str) -> str:
    df: pd.DataFrame = session.get("df")
    cols = list(df.columns) if df is not None else []
    numeric_cols = [c for c in cols if df is not None and pd.api.types.is_numeric_dtype(df[c])]
    cols_json = json.dumps(cols)
    numeric_json = json.dumps(numeric_cols[:4])
    jsx = (
        f"import {{ BarChart, Bar, LineChart, Line, XAxis, YAxis, CartesianGrid, Tooltip, Legend, ResponsiveContainer }} from 'recharts';\n\n"
        f"export default function Dashboard({{ data }}) {{\n"
        f"  // data: array of row objects with keys: {cols_json}\n"
        f"  // numeric columns: {numeric_json}\n"
        f"  return (\n"
        f"    <div className=\"p-6 bg-gray-50 min-h-screen\">\n"
        f"      <h1 className=\"text-2xl font-bold text-blue-900 mb-6\">Analytics Dashboard</h1>\n"
        f"      <div className=\"grid grid-cols-1 md:grid-cols-2 gap-6\">\n"
        f"        <div className=\"bg-white rounded-2xl p-4 shadow\">\n"
        f"          <h2 className=\"font-semibold text-gray-700 mb-3\">{(numeric_cols[:1] or ['value'])[0]} Over Time</h2>\n"
        f"          <ResponsiveContainer width=\"100%\" height={{300}}>\n"
        f"            <LineChart data={{data}}>\n"
        f"              <CartesianGrid strokeDasharray=\"3 3\" />\n"
        f"              <XAxis dataKey=\"{(cols[0] if cols else 'x')}\" />\n"
        f"              <YAxis />\n"
        f"              <Tooltip />\n"
        f"              <Legend />\n"
        + "".join(f'              <Line type=\"monotone\" dataKey=\"{c}\" stroke=\"#2563eb\" dot={{false}} />\n' for c in numeric_cols[:2])
        + f"            </LineChart>\n"
        f"          </ResponsiveContainer>\n"
        f"        </div>\n"
        f"        <div className=\"bg-white rounded-2xl p-4 shadow\">\n"
        f"          <h2 className=\"font-semibold text-gray-700 mb-3\">Breakdown</h2>\n"
        f"          <ResponsiveContainer width=\"100%\" height={{300}}>\n"
        f"            <BarChart data={{data}}>\n"
        f"              <CartesianGrid strokeDasharray=\"3 3\" />\n"
        f"              <XAxis dataKey=\"{(cols[0] if cols else 'x')}\" />\n"
        f"              <YAxis />\n"
        f"              <Tooltip />\n"
        f"              <Legend />\n"
        + "".join(f'              <Bar dataKey=\"{c}\" fill=\"#7c3aed\" />\n' for c in numeric_cols[:2])
        + f"            </BarChart>\n"
        f"          </ResponsiveContainer>\n"
        f"        </div>\n"
        f"      </div>\n"
        f"    </div>\n"
        f"  );\n"
        f"}}"
    )
    safe_jsx = html.escape(jsx)
    return (
        "<div style='background:#1e293b;border-radius:16px;padding:20px;box-shadow:0 1px 6px rgba(0,0,0,.15)'>"
        "<div style='display:flex;justify-content:space-between;align-items:center;margin-bottom:12px'>"
        "<div style='font:700 14px Arial,sans-serif;color:#94a3b8'>&#128196; React JSX Component Scaffold</div>"
        "<button onclick=\"navigator.clipboard.writeText(document.getElementById('jsxCode').innerText)\" style='background:#2563eb;color:#fff;border:none;border-radius:8px;padding:6px 14px;font:600 12px Arial,sans-serif;cursor:pointer'>Copy</button>"
        "</div>"
        f"<pre id='jsxCode' style='font:13px/1.6 monospace;color:#e2e8f0;overflow-x:auto;white-space:pre-wrap'>{safe_jsx}</pre>"
        "</div>"
    )


def _render_dashboard_html(
    session_id: str,
    filename: str,
    question: str,
    analysis: str,
    artifacts: list[dict],
    forecast: Optional[dict],
    session: Optional[dict] = None,
) -> str:
    safe_session_id = html.escape(str(session_id))
    safe_filename = html.escape(str(filename))
    safe_question = html.escape(str(question))

    kpi_html = _build_kpi_cards_html(session) if session else ""

    chartjs_result = _build_chartjs_html(session) if session else ("", "")
    chartjs_html = chartjs_result[0] if isinstance(chartjs_result, tuple) else chartjs_result
    chart_init_js = chartjs_result[1] if isinstance(chartjs_result, tuple) else ""

    analysis_html = _analysis_to_html(analysis)

    artifact_cards = ""
    for artifact in artifacts:
        safe_title = html.escape(str(artifact["title"]))
        safe_url = html.escape(str(artifact["url"]))
        safe_chart_type = html.escape(str(artifact.get("chart_type", "chart")))
        artifact_cards += (
            f"<div style='background:#fff;border:1px solid #e5e7eb;border-radius:16px;padding:16px;box-shadow:0 1px 4px rgba(0,0,0,.06)'>"
            f"<div style='font:700 14px Arial,sans-serif;color:#111827;margin-bottom:3px'>{safe_title}</div>"
            f"<span style='display:inline-block;background:#f3f4f6;color:#6b7280;font-size:11px;font-weight:600;padding:2px 8px;border-radius:20px;margin-bottom:10px;text-transform:uppercase;letter-spacing:.5px'>{safe_chart_type}</span>"
            f"<img src='{safe_url}' alt='{safe_title}' onclick=\"openLightbox('{safe_url}')\" style='width:100%;border-radius:10px;border:1px solid #f3f4f6;display:block;margin-bottom:10px;cursor:zoom-in' loading='lazy' />"
            f"<a href='{safe_url}' target='_blank' rel='noopener noreferrer' style='display:inline-block;background:#111827;color:#fff;text-decoration:none;border-radius:8px;padding:7px 14px;font:600 13px Arial,sans-serif;margin-right:8px'>&#128065; Full Size</a>"
            f"<a href='{safe_url}' download style='display:inline-block;color:#2563eb;text-decoration:none;font:600 13px Arial,sans-serif'>&#8595; Download</a>"
            f"</div>"
        )

    if not artifact_cards:
        artifact_cards = "<p style='color:#9ca3af;font-size:14px'>Chart images are generated for CSV / Excel files.</p>"

    forecast_html = ""
    if forecast is not None:
        safe_furl = html.escape(str(forecast["url"]))
        forecast_html = (
            "<section>"
            "<h2 class='section-title'>&#128200; Forecast &amp; Projections</h2>"
            "<div style='background:#eff6ff;border:1px solid #bfdbfe;border-radius:16px;padding:20px'>"
            f"<p style='font-size:14px;color:#1e40af;margin-bottom:14px'>{html.escape(str(forecast['summary']))}</p>"
            f"<img src='{safe_furl}' alt='Forecast chart' onclick=\"openLightbox('{safe_furl}')\" style='width:100%;border-radius:12px;border:1px solid #bfdbfe;display:block;margin-bottom:14px;cursor:zoom-in' loading='lazy' />"
            f"<div style='display:flex;gap:12px;flex-wrap:wrap'>"
            f"<a href='{safe_furl}' target='_blank' rel='noopener noreferrer' style='display:inline-block;background:#1d4ed8;color:#fff;text-decoration:none;border-radius:8px;padding:7px 14px;font:600 13px Arial,sans-serif'>&#128065; Open Full Size</a>"
            f"<a href='{safe_furl}' download style='display:inline-block;color:#2563eb;text-decoration:none;font:600 13px Arial,sans-serif;padding:7px 0'>&#8595; Download</a>"
            "</div></div></section>"
        )

    jsx_html = _build_jsx_snippet(session, analysis) if session else ""

    kpi_section = (
        f"<div class='card'><h2 class='section-title'>&#127942; Key Metrics</h2>{kpi_html}</div>"
        if kpi_html else ""
    )
    chartjs_section = (
        f"<div class='card'><h2 class='section-title'>&#9728; Interactive Charts</h2>{chartjs_html}</div>"
        if chartjs_html else ""
    )
    jsx_section = (
        f"<div class='card'><h2 class='section-title'>&#9881; JSX Component</h2>{jsx_html}</div>"
        if jsx_html else ""
    )

    return f"""
    <!doctype html>
    <html lang='en'>
    <head>
        <meta charset='utf-8'/>
        <meta name='viewport' content='width=device-width, initial-scale=1'/>
        <title>{safe_filename} \u2014 Analytics Dashboard</title>
        <script src='https://cdn.jsdelivr.net/npm/chart.js@4.4.0/dist/chart.umd.min.js'></script>
        <style>
            *{{box-sizing:border-box;margin:0;padding:0}}
            body{{font-family:'Segoe UI',system-ui,Arial,sans-serif;background:#f0f4f8;color:#1f2937;line-height:1.6}}
            .shell{{max-width:1440px;margin:0 auto;padding:24px;display:grid;gap:22px}}
            .hero{{background:linear-gradient(135deg,#0f172a 0%,#1e3a8a 60%,#1d4ed8 100%);color:#f9fafb;border-radius:20px;padding:28px 32px}}
            .hero h1{{font-size:24px;font-weight:800;letter-spacing:-.3px}}
            .hero .meta{{font-size:13px;opacity:.7;display:flex;gap:20px;flex-wrap:wrap;margin-top:10px}}
            .card{{background:#fff;border:1px solid #e5e7eb;border-radius:18px;padding:22px;box-shadow:0 2px 8px rgba(0,0,0,.06)}}
            .section-title{{font-size:17px;font-weight:700;color:#1e3a8a;margin-bottom:14px;display:flex;align-items:center;gap:8px;padding-bottom:10px;border-bottom:2px solid #e0e7ff}}
            .analysis-body{{font-size:14px;line-height:1.75;color:#1f2937}}
            .analysis-body h3{{margin-top:18px}}
            .analysis-body ul{{margin:6px 0 6px 4px}}
            .charts-grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(440px,1fr));gap:18px}}
            table{{width:100%;border-collapse:collapse;font-size:13px}}
            th{{background:#1e3a8a;color:#fff;padding:10px 14px;text-align:left;font-weight:600}}
            td{{padding:9px 14px;border-bottom:1px solid #f3f4f6}}
            tr:nth-child(even) td{{background:#f8fafc}}
            tr:last-child td{{border-bottom:none}}
        </style>
    </head>
    <body>
    <div class='shell'>
        <div class='hero'>
            <h1>&#128202; Analytics Dashboard</h1>
            <div class='meta'>
                <span>&#128196; {safe_filename}</span>
                <span>&#128273; {safe_session_id[:8]}\u2026</span>
                <span>&#128172; {safe_question}</span>
            </div>
        </div>
        {kpi_section}
        <div class='card'>
            <h2 class='section-title'>&#129504; AI Analysis</h2>
            <div class='analysis-body'>{analysis_html}</div>
        </div>
        {chartjs_section}
        <div class='card'>
            <h2 class='section-title'>&#128444; Chart Artifacts</h2>
            <div class='charts-grid'>{artifact_cards}</div>
        </div>
        {forecast_html}
        {jsx_section}
    </div>
    <div id='lb' onclick='closeLB()' style='display:none;position:fixed;inset:0;background:rgba(0,0,0,.88);z-index:9999;align-items:center;justify-content:center;cursor:zoom-out'>
        <img id='lbImg' src='' style='max-width:93vw;max-height:93vh;border-radius:14px;box-shadow:0 8px 48px rgba(0,0,0,.6);object-fit:contain;cursor:default' onclick='event.stopPropagation()'/>
        <div onclick='closeLB()' style='position:fixed;top:18px;right:26px;color:#fff;font-size:38px;line-height:1;cursor:pointer;font-weight:200;text-shadow:0 2px 8px rgba(0,0,0,.5)'>&times;</div>
        <div onclick='prevImg()' style='position:fixed;left:18px;top:50%;transform:translateY(-50%);color:#fff;font-size:44px;cursor:pointer;user-select:none;opacity:.8'>&#8249;</div>
        <div onclick='nextImg()' style='position:fixed;right:18px;top:50%;transform:translateY(-50%);color:#fff;font-size:44px;cursor:pointer;user-select:none;opacity:.8'>&#8250;</div>
    </div>
    <script>
    var _lbUrls=[], _lbIdx=0;
    function openLightbox(url){{_lbUrls=Array.from(document.querySelectorAll('img[onclick^="openLightbox"]')).map(function(i){{return i.src;}});_lbIdx=_lbUrls.indexOf(url);if(_lbIdx<0){{_lbUrls=[url];_lbIdx=0;}}document.getElementById('lbImg').src=url;document.getElementById('lb').style.display='flex';}}
    function closeLB(){{document.getElementById('lb').style.display='none';}}
    function prevImg(){{_lbIdx=(_lbIdx-1+_lbUrls.length)%_lbUrls.length;document.getElementById('lbImg').src=_lbUrls[_lbIdx];}}
    function nextImg(){{_lbIdx=(_lbIdx+1)%_lbUrls.length;document.getElementById('lbImg').src=_lbUrls[_lbIdx];}}
    document.addEventListener('keydown',function(e){{if(document.getElementById('lb').style.display!=='none'){{if(e.key==='Escape')closeLB();if(e.key==='ArrowLeft')prevImg();if(e.key==='ArrowRight')nextImg();}}}});
    {chart_init_js}
    </script>
    </body>
    </html>
    """


# ── /analyze — non-streaming HTML/JSON dashboard (new endpoint) ───────

@app.post("/analyze")
async def analyze(request: ChatRequest, http_request: Request):
    """Non-streaming analysis endpoint. Returns an HTML dashboard by default,
    or a JSON bundle when response_format='json' or 'bundle'.
    Does not replace /chat (which remains a streaming SSE endpoint)."""
    sid = request.session_id or str(uuid4())
    session = _get_or_create_session(sid, request.system_prompt or ANALYSIS_SYSTEM_PROMPT)

    _trim_messages(session)

    if session.get("df") is None and session.get("cloudinary_url"):
        await asyncio.to_thread(_reload_df_from_cloudinary, session)

    ct = session.get("file_content_type", "")
    if session.get("df_summary"):
        message_with_context = f"[File context]\n{session['df_summary']}\n\n[Question]\n{request.message}"
    elif ct in ("image/jpeg", "image/png", "image/gif", "image/webp") and session.get("file_raw"):
        encoded = base64.standard_b64encode(session["file_raw"]).decode("utf-8")
        message_with_context = [
            {"type": "image", "source": {"type": "base64", "media_type": ct, "data": encoded}},
            {"type": "text", "text": request.message},
        ]
    elif ct == "application/pdf" and session.get("file_raw"):
        encoded = base64.standard_b64encode(session["file_raw"]).decode("utf-8")
        message_with_context = [
            {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": encoded}},
            {"type": "text", "text": request.message},
        ]
    elif session.get("file_summary"):
        message_with_context = f"[File context]\n{session['file_summary']}\n\n[Question]\n{request.message}"
    else:
        message_with_context = request.message

    session["messages"].append({"role": "user", "content": message_with_context})
    session["display"].append({"role": "user", "content": request.message})

    requested_format = (request.response_format or "").lower()
    base_url = str(http_request.base_url).rstrip("/")
    wants_json = requested_format in {"json", "bundle"}

    response = await client.messages.create(
        model=chat_model,
        max_tokens=4096,
        system=session["system"],
        messages=session["messages"],
    )
    analysis_text = _extract_assistant_text(response)
    session["messages"].append({"role": "assistant", "content": analysis_text})
    session["display"].append({"role": "assistant", "content": analysis_text})

    artifacts: list = []
    forecast = None
    visual_bundle: dict = {}
    if session.get("df") is not None:
        try:
            visual_bundle = await _build_chat_visual_bundle(sid, session, base_url)
            artifacts = visual_bundle.get("artifacts", [])
            forecast = visual_bundle.get("forecast")
        except Exception:
            pass

    sessions[sid]["last_accessed"] = time()

    if wants_json:
        return {
            "session_id": sid,
            "file": session.get("file"),
            "analysis": analysis_text,
            "visuals": visual_bundle or None,
        }

    html_content = _render_dashboard_html(
        session_id=sid,
        filename=session.get("file") or "Session Chat",
        question=request.message,
        analysis=analysis_text,
        artifacts=artifacts,
        forecast=forecast,
        session=session,
    )
    return HTMLResponse(content=html_content)


# ── /upload-file — file-only storage, no AI call (new endpoint) ───────

@app.post("/upload-file")
async def upload_file_only(
    file: UploadFile = File(...),
    system_prompt: Optional[str] = Form(None),
    session_id: Optional[str] = Form(None),
):
    """Upload a file into a session without triggering an AI call.
    Use /analyze or /chat with the returned session_id to query it.
    Does not replace /upload (which remains a streaming AI-analysis endpoint)."""
    content_type = _normalized_content_type(file)
    raw = await file.read()

    if not raw:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    parsed_df = None
    summary = None

    if content_type == "text/csv":
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="CSV file could not be decoded as UTF-8")
        parsed_df = await asyncio.to_thread(pd.read_csv, io.StringIO(text), low_memory=False)
        summary = await asyncio.to_thread(_df_summary, parsed_df, file.filename)
    elif content_type in EXCEL_TYPES:
        try:
            parsed_df = await asyncio.to_thread(pd.read_excel, io.BytesIO(raw), engine="openpyxl")
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Could not read Excel file: {e}")
        summary = await asyncio.to_thread(_df_summary, parsed_df, file.filename)
    elif content_type in WORD_TYPES:
        try:
            text_content = await asyncio.to_thread(_extract_docx, raw)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Could not read Word file: {e}")
        summary = f"Word document: {file.filename}\n\n{text_content[:2000]}"
    elif content_type in PPTX_TYPES:
        try:
            text_content = await asyncio.to_thread(_extract_pptx, raw)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Could not read PowerPoint file: {e}")
        summary = f"PowerPoint: {file.filename}\n\n{text_content[:2000]}"
    elif content_type.startswith("text/"):
        try:
            text_content = raw.decode("utf-8")
        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="File could not be decoded as UTF-8 text")
        summary = f"Text file: {file.filename}\n\n{text_content[:2000]}"
    elif content_type in ("image/jpeg", "image/png", "image/gif", "image/webp"):
        summary = f"Image file: {file.filename} ({content_type})"
    elif content_type == "application/pdf":
        summary = f"PDF file: {file.filename}"
    else:
        raise HTTPException(
            status_code=415,
            detail=f"Unsupported file type '{content_type}'. Supported: CSV, Excel, DOCX, PPTX, images, PDF, plain text.",
        )

    sid = session_id or str(uuid4())
    session = _get_or_create_session(sid, system_prompt or ANALYSIS_SYSTEM_PROMPT)

    session["file"] = file.filename
    session["file_content_type"] = content_type
    session["file_summary"] = summary
    if parsed_df is not None:
        session["df"] = parsed_df
        session["df_summary"] = summary
    if content_type not in ("text/csv",) and content_type not in EXCEL_TYPES:
        session["file_raw"] = raw
    sessions[sid]["last_accessed"] = time()

    async def _background_upload():
        resource_type = "image" if content_type.startswith("image/") else "raw"
        try:
            result = await asyncio.to_thread(
                cloudinary.uploader.upload,
                io.BytesIO(raw),
                public_id=f"sessions/{sid}/{file.filename}",
                resource_type=resource_type,
                overwrite=True,
            )
            session["cloudinary_id"] = result["public_id"]
            session["cloudinary_url"] = result["secure_url"]
        except Exception as e:
            print(f"Cloudinary upload failed: {e}")
        await _save_state_to_cloudinary(sid, session)

    asyncio.create_task(_background_upload())

    return {
        "session_id": sid,
        "file": file.filename,
        "content_type": content_type,
        "size_bytes": len(raw),
        "rows": len(parsed_df) if parsed_df is not None else None,
        "columns": list(parsed_df.columns) if parsed_df is not None else None,
        "summary_preview": (summary or "")[:300],
        "status": "ready",
        "next": f"POST /analyze  {{\"session_id\": \"{sid}\", \"message\": \"your question here\"}}",
    }


# ── /upload-html — upload + AI + HTML dashboard (new endpoint) ────────

@app.post("/upload-html")
async def upload_and_render_html(
    file: UploadFile = File(...),
    question: str = Form("Create clickable visual artifacts and a forecast from this file"),
    system_prompt: Optional[str] = Form(None),
    session_id: Optional[str] = Form(None),
):
    content_type = _normalized_content_type(file)
    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    sid = session_id or str(uuid4())
    session = _get_or_create_session(
        sid,
        system_prompt or "You are a helpful assistant. Analyse the provided file and answer the user's question precisely.",
    )
    _trim_messages(session)

    parsed_df = None
    summary = None

    if content_type == "text/csv":
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="CSV file could not be decoded as UTF-8")
        parsed_df = await asyncio.to_thread(pd.read_csv, io.StringIO(text), low_memory=False)
        summary = await asyncio.to_thread(_df_summary, parsed_df, file.filename)
        user_content = [{"type": "text", "text": f"{summary}\n\nQuestion: {question}"}]
    elif content_type in EXCEL_TYPES:
        try:
            parsed_df = await asyncio.to_thread(pd.read_excel, io.BytesIO(raw), engine="openpyxl")
        except Exception as exc:
            raise HTTPException(status_code=400, detail=f"Could not read Excel file: {exc}")
        summary = await asyncio.to_thread(_df_summary, parsed_df, file.filename)
        user_content = [{"type": "text", "text": f"{summary}\n\nQuestion: {question}"}]
    elif content_type in WORD_TYPES:
        text_content = await asyncio.to_thread(_extract_docx, raw)
        user_content = [{"type": "text", "text": f"File: {file.filename}\n\n{text_content}\n\n{question}"}]
    elif content_type in PPTX_TYPES:
        text_content = await asyncio.to_thread(_extract_pptx, raw)
        user_content = [{"type": "text", "text": f"File: {file.filename}\n\n{text_content}\n\n{question}"}]
    elif content_type.startswith("text/"):
        try:
            text_content = raw.decode("utf-8")
        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="File could not be decoded as UTF-8 text")
        user_content = [{"type": "text", "text": f"File: {file.filename}\n\n{text_content}\n\n{question}"}]
    elif content_type in ("image/jpeg", "image/png", "image/gif", "image/webp"):
        encoded = base64.standard_b64encode(raw).decode("utf-8")
        user_content = [
            {"type": "image", "source": {"type": "base64", "media_type": content_type, "data": encoded}},
            {"type": "text", "text": question},
        ]
    elif content_type == "application/pdf":
        encoded = base64.standard_b64encode(raw).decode("utf-8")
        user_content = [
            {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": encoded}},
            {"type": "text", "text": question},
        ]
    else:
        raise HTTPException(
            status_code=415,
            detail=f"Unsupported file type '{content_type}'. Supported: CSV, Excel, images (jpeg/png/gif/webp), PDF, and plain text.",
        )

    session["file"] = file.filename
    if parsed_df is not None:
        session["df"] = parsed_df
        session["df_summary"] = summary

    session["messages"].append({"role": "user", "content": user_content})
    session["display"].append({"role": "user", "content": question})

    response = await client.messages.create(
        model=model,
        max_tokens=1200,
        system=session["system"],
        messages=session["messages"],
    )
    analysis_text = _extract_assistant_text(response)
    session["messages"].append({"role": "assistant", "content": analysis_text})
    session["display"].append({"role": "assistant", "content": analysis_text})

    artifacts = []
    forecast = None
    try:
        artifacts = await _build_visual_artifacts(sid, session, chart_type="auto", max_charts=3)
        forecast = await _build_forecast_artifact(sid, session)
    except HTTPException:
        artifacts = []
        forecast = None

    async def _background_upload():
        resource_type = "image" if content_type.startswith("image/") else "raw"
        try:
            result = await asyncio.to_thread(
                cloudinary.uploader.upload,
                io.BytesIO(raw),
                public_id=f"sessions/{sid}/{file.filename}",
                resource_type=resource_type,
                overwrite=True,
            )
            session["cloudinary_id"] = result["public_id"]
            session["cloudinary_url"] = result["secure_url"]
        except Exception as exc:
            print(f"Cloudinary file upload failed: {exc}")
        await _save_state_to_cloudinary(sid, session)

    asyncio.create_task(_background_upload())
    sessions[sid]["last_accessed"] = time()

    html_content = _render_dashboard_html(
        session_id=sid,
        filename=file.filename,
        question=question,
        analysis=analysis_text,
        artifacts=artifacts,
        forecast=forecast,
        session=session,
    )
    return HTMLResponse(content=html_content)


@app.get("/upload-html")
def upload_html_form():
    html_content = """
    <!doctype html>
    <html>
      <head>
        <meta charset='utf-8'/>
        <meta name='viewport' content='width=device-width, initial-scale=1'/>
        <title>Upload HTML Dashboard</title>
      </head>
      <body style='margin:0;background:#f3f4f6;padding:24px;font-family:Arial,sans-serif'>
        <div style='max-width:720px;margin:0 auto;background:#fff;border:1px solid #e5e7eb;border-radius:12px;padding:16px'>
          <h2 style='margin:0 0 10px 0;color:#111827'>Upload File for Visual Artifacts and Forecast</h2>
          <form action='/upload-html' method='post' enctype='multipart/form-data'>
            <label style='display:block;margin-bottom:8px;color:#374151'>File</label>
            <input type='file' name='file' required style='margin-bottom:12px;width:100%' />
            <label style='display:block;margin-bottom:8px;color:#374151'>Question</label>
            <textarea name='question' rows='4' style='width:100%;margin-bottom:12px'>Create clickable visual artifacts and a forecast from this data</textarea>
            <label style='display:block;margin-bottom:8px;color:#374151'>Session ID (optional)</label>
            <input type='text' name='session_id' style='width:100%;margin-bottom:12px' />
            <button type='submit' style='background:#111827;color:#fff;border:none;border-radius:8px;padding:10px 14px;cursor:pointer'>Render HTML Dashboard</button>
          </form>
        </div>
      </body>
    </html>
    """
    return HTMLResponse(content=html_content)


# ── Visualization & forecast endpoints (new) ──────────────────────────

@app.get("/visualize/{session_id}")
async def visualize_session_data(session_id: str, chart_type: str = "auto", max_charts: int = 3):
    session = sessions.get(session_id)
    if session is None:
        session = _get_or_create_session(session_id, "You are a helpful assistant.")

    artifacts = await _build_visual_artifacts(session_id, session, chart_type=chart_type, max_charts=max_charts)
    sessions[session_id]["last_accessed"] = time()
    return {"session_id": session_id, "count": len(artifacts), "artifacts": artifacts}


@app.get("/forecast/{session_id}")
async def forecast_session_data(session_id: str, periods: int = 3):
    session = sessions.get(session_id)
    if session is None:
        session = _get_or_create_session(session_id, "You are a helpful assistant.")

    forecast = await _build_forecast_artifact(session_id, session, periods=periods)
    sessions[session_id]["last_accessed"] = time()
    return {"session_id": session_id, "forecast": forecast}


@app.get("/artifacts/{session_id}")
async def list_session_artifacts(session_id: str):
    try:
        def _fetch():
            resources = []
            prefix = f"sessions/{session_id}/artifacts/"
            next_cursor = None
            while True:
                kwargs = {
                    "resource_type": "image",
                    "type": "upload",
                    "prefix": prefix,
                    "max_results": 500,
                }
                if next_cursor:
                    kwargs["next_cursor"] = next_cursor
                response = cloudinary.api.resources(**kwargs)
                resources.extend(response.get("resources", []))
                next_cursor = response.get("next_cursor")
                if not next_cursor:
                    break
            return resources

        resources = await asyncio.to_thread(_fetch)
        artifacts = [
            {
                "public_id": resource["public_id"],
                "url": resource["secure_url"],
                "bytes": resource.get("bytes"),
                "format": resource.get("format") or os.path.splitext(resource["public_id"])[-1].lstrip(".") or "unknown",
                "created_at": resource.get("created_at"),
            }
            for resource in resources
        ]
        return {"session_id": session_id, "count": len(artifacts), "artifacts": artifacts}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Cloudinary error: {exc}")


@app.get("/dashboard/{session_id}/render")
async def render_dashboard(session_id: str):
    session = sessions.get(session_id)
    if session is None:
        session = _get_or_create_session(session_id, "You are a helpful assistant.")

    try:
        artifacts = await _build_visual_artifacts(session_id, session, chart_type="auto", max_charts=3)
        forecast = await _build_forecast_artifact(session_id, session)
    except HTTPException:
        artifacts = []
        forecast = None

    html_content = _render_dashboard_html(
        session_id=session_id,
        filename=session.get("file") or "Session Data",
        question="Open dashboard",
        analysis="Dashboard generated from the uploaded session data.",
        artifacts=artifacts,
        forecast=forecast,
        session=session,
    )
    sessions[session_id]["last_accessed"] = time()
    return HTMLResponse(content=html_content)


@app.get("/visualize/{session_id}/render")
async def render_first_artifact(session_id: str):
    session = sessions.get(session_id)
    if session is None:
        session = _get_or_create_session(session_id, "You are a helpful assistant.")
    artifacts = await _build_visual_artifacts(session_id, session, chart_type="auto", max_charts=1)
    if not artifacts:
        raise HTTPException(status_code=404, detail="No artifacts generated for this session.")
    resp = requests.get(artifacts[0]["url"], timeout=30)
    resp.raise_for_status()
    return Response(content=resp.content, media_type="image/png")


# ── /allChat ─ full dashboard, JSX scaffold, shareable links ─────────

class AllChatRequest(BaseModel):
    message: str
    session_id: Optional[str] = None
    system_prompt: Optional[str] = None
    response_format: Optional[str] = None  # "html" (default) | "json"


def _build_allchat_html(
    sid: str,
    filename: str,
    question: str,
    analysis: str,
    artifacts: list[dict],
    forecast: Optional[dict],
    session: Optional[dict],
    base_url: str,
) -> str:
    """Render a richly featured single-page HTML response for /allChat.

    Features:
    • KPI metric cards
    • Interactive Chart.js charts (inline – no external image requests)
    • Cloudinary PNG artifact cards with one-click lightbox
    • Forecast chart card + projected rows table
    • Copyable React/JSX scaffold for any front-end
    • Shareable link panel (dashboard, artifacts, forecast JSON endpoints)
    """
    import re

    safe_sid = html.escape(sid)
    safe_file = html.escape(filename)
    safe_q = html.escape(question)

    # ── KPI cards ──────────────────────────────────────────────────────
    kpi_html = _build_kpi_cards_html(session) if session else ""

    # ── Chart.js interactive visuals ───────────────────────────────────
    chartjs_result = _build_chartjs_html(session) if session else ("", "")
    cjs_html = chartjs_result[0] if isinstance(chartjs_result, tuple) else chartjs_result
    cjs_init = chartjs_result[1] if isinstance(chartjs_result, tuple) else ""

    # ── AI analysis → styled HTML ──────────────────────────────────────
    analysis_html = _analysis_to_html(analysis)

    # ── Artifact cards ─────────────────────────────────────────────────
    artifact_cards = ""
    for a in artifacts:
        t = html.escape(str(a["title"]))
        u = html.escape(str(a["url"]))
        ct = html.escape(str(a.get("chart_type", "chart")))
        artifact_cards += f"""
<div class='ac-card'>
  <div class='ac-card-title'>{t}</div>
  <span class='ac-badge'>{ct}</span>
  <a href='{u}' target='_blank' rel='noopener noreferrer' class='ac-img-link' title='Click to open full-size in browser'>
    <img src='{u}' alt='{t}' loading='lazy' class='ac-img' onclick="openLB('{u}');event.preventDefault()"/>
    <div class='ac-img-overlay'>&#128269; Open</div>
  </a>
  <div class='ac-actions'>
    <a href='{u}' target='_blank' rel='noopener noreferrer' class='btn-dark'>&#128065; Full Screen</a>
    <a href='{u}' download class='btn-link'>&#8595; Download PNG</a>
  </div>
</div>"""

    if not artifact_cards:
        artifact_cards = "<p class='muted'>Upload a CSV or Excel file to generate chart images.</p>"

    # ── Forecast card ──────────────────────────────────────────────────
    forecast_html = ""
    if forecast:
        fu = html.escape(str(forecast["url"]))
        fsummary = html.escape(str(forecast.get("summary", "")))
        rows = forecast.get("forecast_rows", [])
        x_key = forecast.get("x", "date")
        y_key = forecast.get("y", "value")
        table_rows = "".join(
            f"<tr><td>{html.escape(str(r.get(x_key,'')))} </td>"
            f"<td style='text-align:right'>{float(r.get(y_key, 0)):,.2f}</td></tr>"
            for r in rows
        )
        forecast_html = f"""
<section class='card'>
  <h2 class='section-title'>&#128200; Forecast &amp; Projections</h2>
  <p class='forecast-summary'>{fsummary}</p>
  <a href='{fu}' target='_blank' rel='noopener noreferrer' class='ac-img-link' title='Open forecast chart'>
    <img src='{fu}' alt='Forecast' loading='lazy' class='ac-img' onclick="openLB('{fu}');event.preventDefault()"/>
    <div class='ac-img-overlay'>&#128269; Open</div>
  </a>
  <div class='ac-actions' style='margin-top:10px'>
    <a href='{fu}' target='_blank' rel='noopener noreferrer' class='btn-dark'>&#128065; Full Screen</a>
    <a href='{fu}' download class='btn-link'>&#8595; Download PNG</a>
  </div>
  {f"<table class='fc-table'><thead><tr><th>{html.escape(x_key)}</th><th style='text-align:right'>Projected {html.escape(y_key)}</th></tr></thead><tbody>{table_rows}</tbody></table>" if table_rows else ""}
</section>"""

    # ── JSX scaffold ───────────────────────────────────────────────────
    jsx_html = _build_jsx_snippet(session, analysis) if session else ""

    # ── Link panel ─────────────────────────────────────────────────────
    dashboard_url   = f"{base_url}/dashboard/{sid}/render"
    artifacts_url   = f"{base_url}/artifacts/{sid}"
    forecast_url    = f"{base_url}/forecast/{sid}"
    visualize_url   = f"{base_url}/visualize/{sid}"
    safe_du = html.escape(dashboard_url)
    safe_au = html.escape(artifacts_url)
    safe_fu = html.escape(forecast_url)
    safe_vu = html.escape(visualize_url)

    link_panel = f"""
<section class='card'>
  <h2 class='section-title'>&#128279; Shareable Links</h2>
  <div class='link-grid'>
    <div class='link-item'>
      <div class='link-label'>&#127968; Full Dashboard</div>
      <a href='{safe_du}' target='_blank' rel='noopener noreferrer' class='link-url'>{safe_du}</a>
      <button class='copy-btn' onclick="copyLink('{safe_du}',this)">Copy</button>
    </div>
    <div class='link-item'>
      <div class='link-label'>&#128444; Chart Artifacts JSON</div>
      <a href='{safe_au}' target='_blank' rel='noopener noreferrer' class='link-url'>{safe_au}</a>
      <button class='copy-btn' onclick="copyLink('{safe_au}',this)">Copy</button>
    </div>
    <div class='link-item'>
      <div class='link-label'>&#128200; Forecast JSON</div>
      <a href='{safe_fu}' target='_blank' rel='noopener noreferrer' class='link-url'>{safe_fu}</a>
      <button class='copy-btn' onclick="copyLink('{safe_fu}',this)">Copy</button>
    </div>
    <div class='link-item'>
      <div class='link-label'>&#9889; Visualize JSON</div>
      <a href='{safe_vu}' target='_blank' rel='noopener noreferrer' class='link-url'>{safe_vu}</a>
      <button class='copy-btn' onclick="copyLink('{safe_vu}',this)">Copy</button>
    </div>
  </div>
</section>"""

    # ── Assemble page ──────────────────────────────────────────────────
    return f"""<!doctype html>
<html lang='en'>
<head>
<meta charset='utf-8'/>
<meta name='viewport' content='width=device-width,initial-scale=1'/>
<title>{safe_file} — allChat Dashboard</title>
<script src='https://cdn.jsdelivr.net/npm/chart.js@4.4.0/dist/chart.umd.min.js'></script>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:'Segoe UI',system-ui,Arial,sans-serif;background:#f0f4f8;color:#1f2937;line-height:1.6}}
.shell{{max-width:1440px;margin:0 auto;padding:24px;display:grid;gap:22px}}
/* Hero */
.hero{{background:linear-gradient(135deg,#0f172a 0%,#1e3a8a 55%,#1d4ed8 100%);color:#f9fafb;border-radius:20px;padding:28px 32px}}
.hero h1{{font-size:26px;font-weight:800;letter-spacing:-.4px}}
.hero .meta{{font-size:13px;opacity:.72;display:flex;gap:18px;flex-wrap:wrap;margin-top:10px}}
.hero .meta span{{background:rgba(255,255,255,.1);padding:3px 10px;border-radius:20px}}
/* Cards */
.card{{background:#fff;border:1px solid #e5e7eb;border-radius:18px;padding:22px;box-shadow:0 2px 8px rgba(0,0,0,.06)}}
.section-title{{font-size:17px;font-weight:700;color:#1e3a8a;margin-bottom:16px;padding-bottom:10px;border-bottom:2px solid #e0e7ff;display:flex;align-items:center;gap:8px}}
/* Analysis */
.analysis-body{{font-size:14px;line-height:1.8;color:#1f2937}}
.analysis-body h3{{margin-top:18px;font-size:15px;font-weight:700;color:#1e3a8a}}
.analysis-body ul{{margin:6px 0 6px 4px}}
.analysis-body li{{margin:4px 0 4px 18px}}
/* Artifact cards */
.charts-grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(420px,1fr));gap:18px}}
.ac-card{{background:#fff;border:1px solid #e5e7eb;border-radius:16px;padding:16px;box-shadow:0 1px 4px rgba(0,0,0,.06);display:flex;flex-direction:column;gap:8px}}
.ac-card-title{{font:700 14px Arial,sans-serif;color:#111827}}
.ac-badge{{display:inline-block;background:#f3f4f6;color:#6b7280;font:600 11px Arial,sans-serif;padding:2px 9px;border-radius:20px;text-transform:uppercase;letter-spacing:.5px;width:fit-content}}
.ac-img-link{{display:block;position:relative;border-radius:10px;overflow:hidden;cursor:pointer;border:1px solid #f3f4f6}}
.ac-img{{width:100%;display:block;transition:transform .25s}}
.ac-img-link:hover .ac-img{{transform:scale(1.02)}}
.ac-img-overlay{{position:absolute;inset:0;background:rgba(0,0,0,.35);color:#fff;font:700 16px Arial,sans-serif;display:flex;align-items:center;justify-content:center;opacity:0;transition:opacity .2s;letter-spacing:.4px}}
.ac-img-link:hover .ac-img-overlay{{opacity:1}}
.ac-actions{{display:flex;gap:10px;align-items:center;flex-wrap:wrap}}
/* Buttons */
.btn-dark{{display:inline-block;background:#111827;color:#fff;text-decoration:none;border-radius:8px;padding:7px 14px;font:600 13px Arial,sans-serif}}
.btn-dark:hover{{background:#1e3a8a}}
.btn-link{{display:inline-block;color:#2563eb;text-decoration:none;font:600 13px Arial,sans-serif}}
.btn-link:hover{{text-decoration:underline}}
/* Forecast table */
.fc-table{{width:100%;border-collapse:collapse;font-size:13px;margin-top:14px}}
.fc-table th{{background:#1e3a8a;color:#fff;padding:9px 14px;text-align:left;font-weight:600}}
.fc-table td{{padding:8px 14px;border-bottom:1px solid #f3f4f6}}
.fc-table tr:last-child td{{border-bottom:none}}
.fc-table tr:nth-child(even) td{{background:#f8fafc}}
/* Forecast summary */
.forecast-summary{{font-size:14px;color:#1e40af;background:#eff6ff;border:1px solid #bfdbfe;border-radius:10px;padding:12px 16px;margin-bottom:14px}}
/* Link panel */
.link-grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(300px,1fr));gap:14px}}
.link-item{{background:#f8fafc;border:1px solid #e5e7eb;border-radius:12px;padding:14px 16px;display:flex;flex-direction:column;gap:6px}}
.link-label{{font:700 12px Arial,sans-serif;color:#374151;text-transform:uppercase;letter-spacing:.5px}}
.link-url{{font:13px monospace;color:#2563eb;word-break:break-all;text-decoration:none}}
.link-url:hover{{text-decoration:underline}}
.copy-btn{{width:fit-content;background:#e0e7ff;color:#1e3a8a;border:none;border-radius:7px;padding:5px 12px;font:600 12px Arial,sans-serif;cursor:pointer;transition:background .15s}}
.copy-btn:hover{{background:#c7d2fe}}
.copy-btn.copied{{background:#d1fae5;color:#065f46}}
/* Misc */
.muted{{color:#9ca3af;font-size:14px}}
/* Lightbox */
#lb{{display:none;position:fixed;inset:0;background:rgba(0,0,0,.9);z-index:9999;align-items:center;justify-content:center;cursor:zoom-out}}
#lbImg{{max-width:92vw;max-height:92vh;border-radius:12px;box-shadow:0 8px 48px rgba(0,0,0,.6);object-fit:contain;cursor:default}}
#lbClose{{position:fixed;top:16px;right:24px;color:#fff;font-size:40px;line-height:1;cursor:pointer;font-weight:200;opacity:.8}}
#lbPrev,#lbNext{{position:fixed;top:50%;transform:translateY(-50%);color:#fff;font-size:50px;cursor:pointer;user-select:none;opacity:.7;padding:0 10px}}
#lbPrev{{left:8px}}#lbNext{{right:8px}}
</style>
</head>
<body>
<div class='shell'>

<!-- Hero -->
<div class='hero'>
  <h1>&#128202; allChat Dashboard</h1>
  <div class='meta'>
    <span>&#128196; {safe_file}</span>
    <span>&#128273; {safe_sid[:8]}&#8230;</span>
    <span>&#128172; {safe_q}</span>
  </div>
</div>

<!-- KPI Cards -->
{f"<div class='card'><h2 class='section-title'>&#127942; Key Metrics</h2>{kpi_html}</div>" if kpi_html else ""}

<!-- AI Analysis -->
<div class='card'>
  <h2 class='section-title'>&#129504; AI Analysis</h2>
  <div class='analysis-body'>{analysis_html}</div>
</div>

<!-- Interactive Charts -->
{f"<div class='card'><h2 class='section-title'>&#9728;&#65039; Interactive Charts</h2>{cjs_html}</div>" if cjs_html else ""}

<!-- Artifact Cards (Cloudinary PNGs) -->
<div class='card'>
  <h2 class='section-title'>&#128444;&#65039; Chart Artifacts — click any chart to open full size</h2>
  <div class='charts-grid'>{artifact_cards}</div>
</div>

<!-- Forecast -->
{forecast_html}

<!-- JSX Component -->
{f"<div class='card'><h2 class='section-title'>&#9881;&#65039; React / JSX Scaffold — copy &amp; drop into any project</h2>{jsx_html}</div>" if jsx_html else ""}

<!-- Shareable Links -->
{link_panel}

</div><!-- /shell -->

<!-- Lightbox -->
<div id='lb' onclick='closeLB()'>
  <span id='lbClose' onclick='closeLB()'>&#215;</span>
  <span id='lbPrev' onclick='prevImg();event.stopPropagation()'>&#8249;</span>
  <img id='lbImg' src='' onclick='event.stopPropagation()'/>
  <span id='lbNext' onclick='nextImg();event.stopPropagation()'>&#8250;</span>
</div>

<script>
// ── Lightbox ──────────────────────────────────────────────────────────
var _lbUrls=[], _lbIdx=0;
function _gatherImgs(){{return Array.from(document.querySelectorAll('.ac-img')).map(function(i){{return i.src;}});}}
function openLB(url){{
  _lbUrls=_gatherImgs();
  _lbIdx=_lbUrls.indexOf(url);
  if(_lbIdx<0){{_lbUrls=[url];_lbIdx=0;}}
  document.getElementById('lbImg').src=url;
  var lb=document.getElementById('lb');lb.style.display='flex';
}}
function closeLB(){{document.getElementById('lb').style.display='none';}}
function prevImg(){{_lbIdx=(_lbIdx-1+_lbUrls.length)%_lbUrls.length;document.getElementById('lbImg').src=_lbUrls[_lbIdx];}}
function nextImg(){{_lbIdx=(_lbIdx+1)%_lbUrls.length;document.getElementById('lbImg').src=_lbUrls[_lbIdx];}}
document.addEventListener('keydown',function(e){{
  if(document.getElementById('lb').style.display!=='none'){{
    if(e.key==='Escape')closeLB();
    if(e.key==='ArrowLeft')prevImg();
    if(e.key==='ArrowRight')nextImg();
  }}
}});

// ── Copy link button ──────────────────────────────────────────────────
function copyLink(url,btn){{
  navigator.clipboard.writeText(url).then(function(){{
    btn.textContent='Copied!';btn.classList.add('copied');
    setTimeout(function(){{btn.textContent='Copy';btn.classList.remove('copied');}},2000);
  }});
}}

// ── Chart.js init ─────────────────────────────────────────────────────
{cjs_init}
</script>
</body>
</html>"""


@app.post("/allChat")
async def all_chat(request: AllChatRequest, http_request: Request):
    """
    All-in-one endpoint: runs AI analysis, generates interactive Chart.js charts,
    Cloudinary PNG artifact cards (click-to-open), a forecast, a copyable JSX
    component, and a shareable link panel — all in a single HTML response.

    Pass response_format='json' to get a structured JSON bundle instead.
    """
    sid = request.session_id or str(uuid4())
    session = _get_or_create_session(sid, request.system_prompt or ANALYSIS_SYSTEM_PROMPT)
    _trim_messages(session)

    if session.get("df") is None and session.get("cloudinary_url"):
        await asyncio.to_thread(_reload_df_from_cloudinary, session)

    # Build user message with any file context
    ct = session.get("file_content_type", "")
    if session.get("df_summary"):
        msg_content = f"[File context]\n{session['df_summary']}\n\n[Question]\n{request.message}"
    elif ct in ("image/jpeg", "image/png", "image/gif", "image/webp") and session.get("file_raw"):
        enc = base64.standard_b64encode(session["file_raw"]).decode("utf-8")
        msg_content = [
            {"type": "image", "source": {"type": "base64", "media_type": ct, "data": enc}},
            {"type": "text", "text": request.message},
        ]
    elif ct == "application/pdf" and session.get("file_raw"):
        enc = base64.standard_b64encode(session["file_raw"]).decode("utf-8")
        msg_content = [
            {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": enc}},
            {"type": "text", "text": request.message},
        ]
    elif session.get("file_summary"):
        msg_content = f"[File context]\n{session['file_summary']}\n\n[Question]\n{request.message}"
    else:
        msg_content = request.message

    session["messages"].append({"role": "user", "content": msg_content})
    session["display"].append({"role": "user", "content": request.message})

    ai_response = await client.messages.create(
        model=chat_model,
        max_tokens=4096,
        system=session["system"],
        messages=session["messages"],
    )
    analysis_text = _extract_assistant_text(ai_response)
    session["messages"].append({"role": "assistant", "content": analysis_text})
    session["display"].append({"role": "assistant", "content": analysis_text})

    # Build visuals if tabular data is present
    artifacts: list = []
    forecast: Optional[dict] = None
    if session.get("df") is not None:
        try:
            artifacts = await _build_visual_artifacts(sid, session, chart_type="auto", max_charts=6)
        except Exception:
            artifacts = []
        try:
            forecast = await _build_forecast_artifact(sid, session)
        except Exception:
            forecast = None

    sessions[sid]["last_accessed"] = time()
    asyncio.create_task(_save_state_to_cloudinary(sid, session))

    base_url = str(http_request.base_url).rstrip("/")
    wants_json = (request.response_format or "").lower() in {"json", "bundle"}

    if wants_json:
        return {
            "session_id": sid,
            "file": session.get("file"),
            "analysis": analysis_text,
            "artifacts": artifacts,
            "forecast": forecast,
            "links": {
                "dashboard": f"{base_url}/dashboard/{sid}/render",
                "artifacts": f"{base_url}/artifacts/{sid}",
                "forecast": f"{base_url}/forecast/{sid}",
                "visualize": f"{base_url}/visualize/{sid}",
            },
        }

    html_content = _build_allchat_html(
        sid=sid,
        filename=session.get("file") or "Chat Session",
        question=request.message,
        analysis=analysis_text,
        artifacts=artifacts,
        forecast=forecast,
        session=session,
        base_url=base_url,
    )
    return HTMLResponse(content=html_content)


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="0.0.0.0", port=port, reload=True, log_level="warning")
