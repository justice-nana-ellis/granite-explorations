"""Unified file parsing service.

Accepts a FastAPI UploadFile (or raw bytes + content_type) and returns a
ParsedFile dataclass containing everything routes need to build messages
and populate the session.
"""
import asyncio
import base64
import io
import mimetypes
from dataclasses import dataclass, field
from typing import Optional

import pandas as pd
from docx import Document as DocxDocument
from fastapi import HTTPException, UploadFile
from pptx import Presentation

from app.utils.df_utils import df_summary

EXCEL_TYPES = frozenset({
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
})

WORD_TYPES = frozenset({
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
})

PPTX_TYPES = frozenset({
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
})

IMAGE_TYPES = frozenset({"image/jpeg", "image/png", "image/gif", "image/webp"})


@dataclass
class ParsedFile:
    filename: str
    content_type: str
    raw: bytes
    df: Optional[pd.DataFrame] = None
    df_summary_text: Optional[str] = None
    file_summary: Optional[str] = None
    file_raw: Optional[bytes] = None
    user_content: list = field(default_factory=list)


def normalize_content_type(upload: UploadFile) -> str:
    ct = upload.content_type or ""
    if not ct or "/" not in ct or ct == "application/octet-stream":
        guessed, _ = mimetypes.guess_type(upload.filename or "")
        ct = guessed or "application/octet-stream"
    if ct == "image/jpg":
        ct = "image/jpeg"
    return ct


async def parse_upload(file: UploadFile, question: str) -> ParsedFile:
    content_type = normalize_content_type(file)
    raw = await file.read()

    if not raw:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    return await _parse_raw(raw, file.filename or "file", content_type, question)


async def _parse_raw(raw: bytes, filename: str, content_type: str, question: str) -> ParsedFile:
    result = ParsedFile(filename=filename, content_type=content_type, raw=raw)

    if content_type == "text/csv":
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="CSV file could not be decoded as UTF-8")
        result.df = await asyncio.to_thread(pd.read_csv, io.StringIO(text), low_memory=False)
        result.df_summary_text = await asyncio.to_thread(df_summary, result.df, filename)
        result.user_content = [{"type": "text", "text": f"{result.df_summary_text}\n\nQuestion: {question}"}]

    elif content_type in EXCEL_TYPES:
        try:
            result.df = await asyncio.to_thread(pd.read_excel, io.BytesIO(raw), engine="openpyxl")
        except Exception as exc:
            raise HTTPException(status_code=400, detail=f"Could not read Excel file: {exc}")
        result.df_summary_text = await asyncio.to_thread(df_summary, result.df, filename)
        result.user_content = [{"type": "text", "text": f"{result.df_summary_text}\n\nQuestion: {question}"}]

    elif content_type in WORD_TYPES:
        try:
            text_content = await asyncio.to_thread(_extract_docx, raw)
        except Exception as exc:
            raise HTTPException(status_code=400, detail=f"Could not read Word file: {exc}")
        result.file_summary = f"Word document: {filename}\n\n{text_content[:2000]}"
        result.user_content = [{"type": "text", "text": f"File: {filename}\n\n{text_content}\n\n{question}"}]

    elif content_type in PPTX_TYPES:
        try:
            text_content = await asyncio.to_thread(_extract_pptx, raw)
        except Exception as exc:
            raise HTTPException(status_code=400, detail=f"Could not read PowerPoint file: {exc}")
        result.file_summary = f"PowerPoint: {filename}\n\n{text_content[:2000]}"
        result.user_content = [{"type": "text", "text": f"File: {filename}\n\n{text_content}\n\n{question}"}]

    elif content_type.startswith("text/"):
        try:
            text_content = raw.decode("utf-8")
        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="File could not be decoded as UTF-8 text")
        result.file_summary = f"Text file: {filename}\n\n{text_content[:2000]}"
        result.user_content = [{"type": "text", "text": f"File: {filename}\n\n{text_content}\n\n{question}"}]

    elif content_type in IMAGE_TYPES:
        encoded = base64.standard_b64encode(raw).decode("utf-8")
        result.file_raw = raw
        result.file_summary = f"Image file: {filename} ({content_type})"
        result.user_content = [
            {"type": "image", "source": {"type": "base64", "media_type": content_type, "data": encoded}},
            {"type": "text", "text": question},
        ]

    elif content_type == "application/pdf":
        encoded = base64.standard_b64encode(raw).decode("utf-8")
        result.file_raw = raw
        result.file_summary = f"PDF file: {filename}"
        result.user_content = [
            {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": encoded}},
            {"type": "text", "text": question},
        ]

    else:
        raise HTTPException(
            status_code=415,
            detail=(
                f"Unsupported file type '{content_type}'. "
                "Supported: CSV, Excel, DOCX, PPTX, images (jpeg/png/gif/webp), PDF, plain text."
            ),
        )

    return result


def _extract_docx(raw: bytes) -> str:
    doc = DocxDocument(io.BytesIO(raw))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(raw: bytes) -> str:
    prs = Presentation(io.BytesIO(raw))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)
